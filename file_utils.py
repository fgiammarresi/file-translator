# Definisci una funzione per estrarre il testo da diversi formati di file
def extract_text_from_file(file_path):
    text = ""
    try:
        if file_path.endswith(".pdf"):
            try:
                import PyPDF2
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    for page_num in range(len(reader.pages)):
                        page = reader.pages[page_num]
                        text += page.extract_text()
            except ImportError:
                return "Libreria PyPDF2 non trovata. Installala con: !pip install PyPDF2"
        elif file_path.endswith(".docx"):
            try:
                from docx import Document
                document = Document(file_path)
                for paragraph in document.paragraphs:
                    text += paragraph.text + "\n"
            except ImportError:
                return "Libreria python-docx non trovata. Installala con: !pip install python-docx"
        elif file_path.endswith(".pptx"):
            try:
                import pptx
                prs = pptx.Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text += run.text + " "
                            text += "\n"
            except ImportError:
                return "Libreria python-pptx non trovata. Installala con: !pip install python-pptx"
        elif file_path.endswith(".xlsx"):
            try:
                import openpyxl
                workbook = openpyxl.load_workbook(file_path)
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    for row in sheet.iter_rows():
                        row_text = " ".join(str(cell.value) for cell in row if cell.value is not None)
                        text += row_text + "\n"
            except ImportError:
                return "Libreria openpyxl non trovata. Installala con: !pip install openpyxl"
        else:
            return "Formato file non supportato."
        return text
    except Exception as e:
        return f"Si è verificato un errore durante l'estrazione del testo: {e}"

# Esempio di utilizzo (questo non verrà eseguito direttamente su GitHub)
# Dovremo eseguirlo in Colab dopo aver clonato il repository.
# file_path_colab = '/content/esempio.pdf'
# extracted_text = extract_text_from_file(file_path_colab)
# print(extracted_text[:200] if isinstance(extracted_text, str) else extracted_text)
